"""
Generate VeloxCon 2026 PowerPoint — v2 (post 3-model review).
Incorporates feedback from GPT-5.4, Claude Sonnet 4.5, and Claude Haiku 4.5.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ──────────────────────────────────────────────────────────────────────
# Microsoft Brand Colors
# ──────────────────────────────────────────────────────────────────────
MS_BLUE = RGBColor(0x00, 0x78, 0xD4)
MS_DARK_BLUE = RGBColor(0x00, 0x3D, 0x6B)
MS_LIGHT_BLUE = RGBColor(0xEB, 0xF3, 0xFB)
MS_PURPLE = RGBColor(0x88, 0x6C, 0xE4)
MS_GREEN = RGBColor(0x10, 0x7C, 0x10)
MS_RED = RGBColor(0xD1, 0x34, 0x38)
MS_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
MS_GRAY_10 = RGBColor(0xFA, 0xF9, 0xF8)
MS_GRAY_30 = RGBColor(0xED, 0xEB, 0xE9)
MS_GRAY_60 = RGBColor(0x60, 0x5E, 0x5C)
MS_GRAY_90 = RGBColor(0x32, 0x31, 0x30)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

MS_LOGO_RED = RGBColor(0xF2, 0x50, 0x22)
MS_LOGO_GREEN = RGBColor(0x7F, 0xBA, 0x00)
MS_LOGO_BLUE = RGBColor(0x00, 0xA4, 0xEF)
MS_LOGO_YELLOW = RGBColor(0xFF, 0xB9, 0x00)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ──────────────────────────────────────────────────────────────────────
# Helper Functions
# ──────────────────────────────────────────────────────────────────────
def add_ms_logo(slide, left, top, size=Inches(0.25)):
    gap = Emu(int(size * 0.08))
    sq = Emu(int((size - gap) / 2))
    colors = [MS_LOGO_RED, MS_LOGO_GREEN, MS_LOGO_BLUE, MS_LOGO_YELLOW]
    positions = [(0, 0), (sq + gap, 0), (0, sq + gap), (sq + gap, sq + gap)]
    for (dx, dy), color in zip(positions, colors):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + dx, top + dy, sq, sq)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()


def add_gradient_bg(slide, color1=MS_BLUE, color2=MS_DARK_BLUE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.gradient()
    bg.fill.gradient_stops[0].color.rgb = color1
    bg.fill.gradient_stops[0].position = 0.0
    bg.fill.gradient_stops[1].color.rgb = color2
    bg.fill.gradient_stops[1].position = 1.0
    bg.line.fill.background()
    bg.rotation = 315
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_accent_bar(slide, top=Inches(1.5), color=MS_BLUE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(1.2), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_bottom_bar(slide, text="VeloxCon 2026  |  April 30  |  Meta HQ, Menlo Park"):
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35))
    p = footer.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = MS_GRAY_60
    p.font.name = "Segoe UI"


def add_slide_number(slide, num, color=MS_GRAY_60):
    tb = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{num} / {TOTAL_SLIDES}"
    p.font.size = Pt(9)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = "Segoe UI"


def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=MS_GRAY_90, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Segoe UI"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(2)
    return tf


def add_title_text(slide, left, top, width, height, text, font_size=32, color=MS_GRAY_90):
    return add_textbox(slide, left, top, width, height, text,
                       font_size=font_size, color=color, bold=True,
                       font_name="Segoe UI Semibold")


def add_bullet_list(tf, items, font_size=13, color=MS_GRAY_90):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        p.space_before = Pt(2)
        if ":" in item:
            prefix, rest = item.split(":", 1)
            run = p.add_run()
            run.text = prefix + ":"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = True
            run.font.name = "Segoe UI"
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Segoe UI"
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Segoe UI"


def add_card(slide, left, top, width, height, title, body,
             accent_color=MS_BLUE, bg_color=MS_GRAY_10):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = MS_GRAY_30
    card.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1),
                width - Inches(0.3), Inches(0.35),
                title, font_size=13, color=accent_color, bold=True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.42),
                width - Inches(0.3), height - Inches(0.5),
                body, font_size=11, color=MS_GRAY_60)


def add_metric_card(slide, left, top, width, height, value, label,
                    ref="", accent_color=MS_BLUE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent_color
    card.line.width = Pt(2)
    add_textbox(slide, left, top + Inches(0.12), width, Inches(0.45),
                value, font_size=28, color=accent_color, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.08), top + Inches(0.55),
                width - Inches(0.16), Inches(0.35),
                label, font_size=10, color=MS_GRAY_90,
                alignment=PP_ALIGN.CENTER)
    if ref:
        add_textbox(slide, left + Inches(0.08), top + Inches(0.85),
                    width - Inches(0.16), Inches(0.25),
                    ref, font_size=8, color=MS_GRAY_60,
                    alignment=PP_ALIGN.CENTER)


def add_table(slide, left, top, width, rows, cols, data, col_widths=None,
              header_color=MS_BLUE, font_size=11):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.42 * rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            cell.margin_left = Pt(8)
            cell.margin_right = Pt(8)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Segoe UI"
                paragraph.space_after = Pt(1)
                if r == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = MS_GRAY_90
            cell.text_frame.word_wrap = True
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = MS_GRAY_10
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table_shape


def add_thesis_box(slide, left, top, width, height, prefix, body):
    """Thesis/callout box with blue left bar."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = MS_LIGHT_BLUE
    box.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = MS_BLUE
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.08), width - Inches(0.3), height - Inches(0.15))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = prefix + " "
    run.font.size = Pt(13)
    run.font.color.rgb = MS_BLUE
    run.font.bold = True
    run.font.name = "Segoe UI"
    run2 = p.add_run()
    run2.text = body
    run2.font.size = Pt(13)
    run2.font.color.rgb = MS_GRAY_90
    run2.font.name = "Segoe UI"


def add_speaker_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def section_divider(num, title, subtitle, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide, MS_BLUE, MS_PURPLE)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(2), Inches(0.6),
                f"0{num}", font_size=48, color=WHITE, bold=True,
                font_name="Segoe UI Light")
    add_textbox(slide, Inches(0.8), Inches(2.8), Inches(10), Inches(0.8),
                title, font_size=36, color=WHITE, bold=True,
                font_name="Segoe UI Semibold")
    add_textbox(slide, Inches(0.8), Inches(3.7), Inches(10), Inches(0.5),
                subtitle, font_size=18, color=RGBColor(0xB4, 0xD6, 0xFA))
    add_slide_number(slide, slide_num, color=RGBColor(0xB4, 0xD6, 0xFA))
    return slide


# ======================================================================
# SLIDE 1: TITLE
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, MS_BLUE, MS_DARK_BLUE)
add_ms_logo(slide, Inches(5.9), Inches(1.2), Inches(0.55))
add_textbox(slide, Inches(6.55), Inches(1.22), Inches(2), Inches(0.55),
            "Microsoft", font_size=22, color=WHITE, bold=True)
add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.2),
            "From Apache Incubator to\nTop-Level Project",
            font_size=40, color=WHITE, bold=True,
            font_name="Segoe UI Semibold", alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.5),
            "The Apache Gluten Journey with Velox at Microsoft",
            font_size=22, color=RGBColor(0xB4, 0xD6, 0xFA),
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.5),
            "Production lessons from Fabric NEE — and the upstream contributions they drove in Gluten and Velox",
            font_size=14, color=RGBColor(0x8E, 0xC8, 0xFF),
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.4),
            "William Chen  &  Ajith Shetty  |  Microsoft  |  VeloxCon 2026",
            font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(5.65), Inches(10.3), Inches(0.3),
            "April 30, 2026  |  Meta HQ, Menlo Park",
            font_size=12, color=RGBColor(0x8E, 0xC8, 0xFF),
            alignment=PP_ALIGN.CENTER)


# ======================================================================
# SLIDE 2: AGENDA
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6), "Agenda")
add_accent_bar(slide, top=Inches(0.95))

items = [
    ("01", "  Why This Moment Matters — Gluten TLP & Microsoft's engineering commitment"),
    ("02", "  Architecture in Practice — Where Spark ends, Gluten translates, Velox executes"),
    ("03", "  Production Lessons — What Fabric-scale workloads exposed, what the community upstreamed"),
    ("04", "  Upstream Contributions — Semantic parity, fallback reduction, performance"),
    ("05", "  Fabric NEE — Productionizing Gluten + Velox in Microsoft Fabric"),
    ("06", "  Graduation & Roadmap — TLP lessons & 2026 outlook"),
    ("07", "  Call for Action"),
]
tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11), Inches(3.5))
tf = tb.text_frame
tf.word_wrap = True
for i, (num, rest) in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    run = p.add_run()
    run.text = num
    run.font.size = Pt(15)
    run.font.color.rgb = MS_BLUE
    run.font.bold = True
    run.font.name = "Segoe UI"
    run2 = p.add_run()
    run2.text = rest
    run2.font.size = Pt(15)
    run2.font.color.rgb = MS_GRAY_90
    run2.font.name = "Segoe UI"

# Context note
ctx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.9))
ctx.fill.solid()
ctx.fill.fore_color.rgb = MS_LIGHT_BLUE
ctx.line.fill.background()
add_textbox(slide, Inches(1.0), Inches(5.3), Inches(11.1), Inches(0.7),
            "Context: Yuan Zhou's session (1:40 PM) covers Gluten's community growth and ecosystem broadly. "
            "This talk focuses on production lessons at Fabric scale — and how those lessons "
            "upstream to Gluten, Velox, and Spark.",
            font_size=12, color=MS_GRAY_90)
add_bottom_bar(slide)
add_slide_number(slide, 2)
add_speaker_notes(slide,
    "We're telling the story of Apache Gluten's journey to TLP, centered on production lessons at Fabric scale. "
    "Yuan Zhou covers community growth later today; we focus on how real production users drive upstream innovation. "
    "By the end, you'll understand the architecture, the problems production exposed, and how you can get involved.")


# ======================================================================
# SLIDE 3: SECTION DIVIDER — Why This Moment Matters
# ======================================================================
section_divider(1, "Why This Moment Matters",
                "Production validation as a force multiplier for open-source maturity", 3)


# ======================================================================
# SLIDE 4: Why This Moment Matters — Content
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6), "Why This Moment Matters")
add_accent_bar(slide, top=Inches(0.95))

# Announcement
ann = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.6))
ann.fill.solid()
ann.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
ann.line.color.rgb = MS_GREEN
ann.line.width = Pt(1.5)
add_textbox(slide, Inches(1.0), Inches(1.28), Inches(11.1), Inches(0.4),
            "February 18, 2026 — Apache Gluten graduated to an ASF Top-Level Project.",
            font_size=15, color=MS_GRAY_90, bold=True)

# Thesis box — softened per Sonnet/GPT
add_thesis_box(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.3),
    "Hypothesis:",
    "Production feedback loops accelerate open-source maturity. At Fabric scale, "
    "we discover semantic gaps, edge-case bugs, and performance bottlenecks faster "
    "than synthetic benchmarks — and upstream fixes within sprints. "
    "This production-to-upstream cycle helped Gluten reach TLP graduation.")

# Three role cards
roles = [
    ("Production User", "Fabric NEE at scale", MS_BLUE),
    ("Upstream Contributor", "Gluten + Velox + Spark", MS_GREEN),
    ("Community Partner", "PMC & committer seat", MS_PURPLE),
]
for i, (title, desc, color) in enumerate(roles):
    x = Inches(0.8 + i * 4.0)
    add_metric_card(slide, x, Inches(3.6), Inches(3.6), Inches(1.1),
                    title, desc, accent_color=color)

add_textbox(slide, Inches(0.8), Inches(4.95), Inches(11.5), Inches(0.5),
            "The upstream-first pattern: validate at scale, contribute in the open, "
            "ship products based on the shared foundation.",
            font_size=12, color=MS_GRAY_60, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.25),
            "Ref: incubator.apache.org/projects/gluten.html",
            font_size=9, color=MS_GRAY_60)
add_bottom_bar(slide)
add_slide_number(slide, 4)
add_speaker_notes(slide,
    "February 2026 is significant: Gluten became an ASF Top-Level Project. But here's what made "
    "graduation possible — production deployment at scale. Scale is the fastest feedback loop. "
    "At thousands of concurrent queries, we discover memory bugs, semantic gaps, and performance "
    "bottlenecks that synthetic benchmarks miss. We committed to upstream-first: every fix goes to "
    "Gluten and Velox first, and Fabric productizes on top of shared foundations. "
    "This is not a fork strategy — it's a partnership.")


# ======================================================================
# SLIDE 5: The Choice We Faced
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6), "The Choice We Faced")
add_accent_bar(slide, top=Inches(0.95))

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.5),
            "Fabric runs Spark at large scale. We needed material performance gains "
            "without requiring application rewrites or a separate premium runtime.",
            font_size=14, color=MS_GRAY_90)

table_data = [
    ["Approach", "Assessment"],
    ["Build from scratch", "Years of effort, no community leverage"],
    ["Fork Spark + custom codegen", "Diverges from upstream, maintenance burden"],
    ["Proprietary native engine", "Closed ecosystem, limited upstream leverage"],
    ["Gluten + Velox  \u2713", "Production-proven execution engine, modular integration model, active multi-company community"],
]
tbl = add_table(slide, Inches(0.8), Inches(1.9), Inches(11.5), 5, 2, table_data,
                col_widths=[Inches(3.5), Inches(8.0)], font_size=12)
# Highlight winning row
for c in range(2):
    cell = tbl.table.cell(4, c)
    cell.fill.solid()
    cell.fill.fore_color.rgb = MS_LIGHT_BLUE
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.color.rgb = MS_BLUE

add_thesis_box(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(1.0),
    "Community flywheel:",
    "Production adopters discover issues at scale \u2192 upstream fixes \u2192 "
    "Meta, Microsoft, Intel, Kyligence, and all adopters benefit \u2192 repeat.")

add_bottom_bar(slide)
add_slide_number(slide, 5)
add_speaker_notes(slide,
    "We had four paths. Gluten was already Meta-proven at exabyte scale, with modular backends and "
    "a growing multi-company community. The deciding factor: we could be a production user AND a "
    "contributor. The flywheel on this slide is not aspirational — it's what's actually happening. "
    "Fabric finds a bug, we upstream a fix, the entire community benefits.")


# ======================================================================
# SLIDE 6: SECTION DIVIDER — Architecture
# ======================================================================
section_divider(2, "Architecture in Practice",
                "Where Spark ends, Gluten translates, and Velox takes over", 6)


# ======================================================================
# SLIDE 7: Execution Flow
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
               "Execution Flow: Spark \u2192 Gluten \u2192 Velox")
add_accent_bar(slide, top=Inches(0.95))

# Flow boxes — fixed Substrait terminology per Sonnet
flow_items = [
    ("Spark Planning", "Catalyst optimizer, AQE,\nCBO — unchanged", MS_GRAY_30, MS_GRAY_10),
    ("Gluten Translation", "Catalyst plan \u2192 Substrait IR\n\u2192 Velox operators (via JNI)", MS_BLUE, MS_LIGHT_BLUE),
    ("Velox Execution", "Vectorized C++ | SIMD\nColumnar | Off-heap", MS_GREEN, RGBColor(0xE8, 0xF5, 0xE9)),
]
for i, (title, desc, border, bg) in enumerate(flow_items):
    x = Inches(0.8 + i * 4.2)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(3.5), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = border
    box.line.width = Pt(2)
    add_textbox(slide, x + Inches(0.15), Inches(1.38), Inches(3.2), Inches(0.35),
                title, font_size=14, color=MS_GRAY_90, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), Inches(1.72), Inches(3.2), Inches(0.55),
                desc, font_size=11, color=MS_GRAY_60, alignment=PP_ALIGN.CENTER)
    if i < 2:
        add_textbox(slide, x + Inches(3.5), Inches(1.55), Inches(0.7), Inches(0.5),
                    "\u2192", font_size=28, color=MS_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Key principle — refined per GPT/Sonnet
add_textbox(slide, Inches(0.8), Inches(2.7), Inches(11.5), Inches(0.4),
            "Gluten replaces Spark's row-based Volcano iterator with Velox's vectorized pipeline "
            "— Spark's planning, scheduling, and fault tolerance remain unchanged.",
            font_size=12, color=MS_GRAY_60, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(3.2), Inches(8), Inches(0.4),
            "Simplified Plan with Fallback Boundary", font_size=16, color=MS_GRAY_90, bold=True)

code_text = """== Physical Plan (with Gluten) ==
+- ^(2) ProjectExecTransformer [...]           \u2190 NATIVE: projection
  +- ^(2) ShuffledHashJoinExecTransformer [...]  \u2190 NATIVE: hash join
    :- ColumnarToRowExec                         \u2190 fallback boundary (UnsupportedExpr)
    :  +- ^(1) FilterExecTransformer [...]       \u2190 NATIVE: filter
    :    +- ^(1) NativeFileScan parquet          \u2190 NATIVE: scan with pushdown
    +- ^(0) ColumnarExchange [...]               \u2190 NATIVE: shuffle
      +- ^(0) FilterExecTransformer [...]        \u2190 NATIVE: filter
        +- ^(0) NativeFileScan parquet           \u2190 NATIVE: scan"""

code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.6), Inches(11.5), Inches(2.7))
code_bg.fill.solid()
code_bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
code_bg.line.color.rgb = MS_GRAY_30
code_bg.line.width = Pt(1)
add_textbox(slide, Inches(1.0), Inches(3.7), Inches(11.1), Inches(2.5),
            code_text, font_size=10, color=MS_GRAY_90, font_name="Cascadia Code")

add_textbox(slide, Inches(0.8), Inches(6.45), Inches(11.5), Inches(0.4),
            "Goal: preserve end-to-end columnar execution and minimize ColumnarToRow / RowToColumnar transitions.",
            font_size=12, color=MS_BLUE, bold=True)

add_bottom_bar(slide)
add_slide_number(slide, 7)
add_speaker_notes(slide,
    "Spark's Catalyst optimizer, CBO, AQE — all unchanged. Gluten intercepts the optimized plan, "
    "translates it to Substrait IR, and hands it to Velox via JNI. Velox runs vectorized C++ operators. "
    "The plan example shows a real fallback boundary: ColumnarToRowExec marks where Velox cannot handle "
    "an expression, so data goes back to JVM. Minimizing those transitions is the key optimization goal.")


# ======================================================================
# SLIDE 8: SECTION DIVIDER — Production Lessons
# ======================================================================
section_divider(3, "Production Lessons",
                "Correctness, memory, and observability challenges from real workloads", 8)


# ======================================================================
# SLIDE 9: The Memory Challenge
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
               "Solving Cross-Heap OOMs: A 9-Month Engineering Effort")
add_accent_bar(slide, top=Inches(0.95))

# LEFT: Problem
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.35),
            "The Problem", font_size=18, color=MS_RED, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.8),
            "Early 2025: TPC-DS queries crash with OOM after 90 seconds — despite available memory. "
            "Without shared accounting, Velox's MmapAllocator and JVM heap competed independently "
            "— neither could react to global memory pressure.",
            font_size=12, color=MS_GRAY_90)

add_card(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(2.0),
         "Root Cause & Symptom",
         "No unified memory accounting across JVM on-heap and Velox off-heap (via jemalloc). "
         "Velox allocated from OS while JVM GC ran under pressure, causing container OOMs "
         "before either allocator could react.\n\n"
         "Symptom: Large aggregations (GROUP BY with 100M+ groups) OOM at ~60% reported "
         "memory — Velox off-heap was invisible to Spark's memory manager.",
         accent_color=MS_RED, bg_color=RGBColor(0xFD, 0xF0, 0xF0))

# RIGHT: Fix
add_textbox(slide, Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.35),
            "The Fix & Impact", font_size=18, color=MS_GREEN, bold=True)
add_textbox(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(0.4),
            "9-month collaboration across Gluten, Velox, and Fabric teams:",
            font_size=12, color=MS_GRAY_90)

fix_items = [
    "Reservation protocol: Velox requests quota from Spark's MemoryManager before allocation",
    "GC pressure signals: JVM broadcasts pressure events \u2192 triggers Velox spill before OOM",
    "Graceful degradation: Under contention, Velox spills sorted partitions to SSD with async I/O",
]
tb = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.7), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True
add_bullet_list(tf, fix_items, font_size=11)

add_metric_card(slide, Inches(6.8), Inches(3.7), Inches(2.7), Inches(1.15),
                "\u2193 70%", "OOM incidents\n(Fabric production telemetry)",
                ref="Gluten 1.3.0 + Runtime 1.3", accent_color=MS_GREEN)
add_metric_card(slide, Inches(9.8), Inches(3.7), Inches(2.7), Inches(1.15),
                "\u2192 Upstream", "Gluten 1.3.0 + community",
                ref="Runtime 1.3 Memory Blog", accent_color=MS_BLUE)

add_bottom_bar(slide)
add_slide_number(slide, 9)
add_speaker_notes(slide,
    "This is the story that illustrates why Gluten needed production feedback. January 2025: TPC-DS queries "
    "crashed with OOM despite available memory. The root cause was subtle: Velox's MmapAllocator and "
    "JVM GC competed independently for the same container memory. We spent 9 months building a cooperative "
    "memory arbitration protocol. The result: OOM incidents dropped 70% in Fabric production telemetry. "
    "This fix went upstream in Gluten 1.3.0 and now benefits every Gluten user.")


# ======================================================================
# SLIDE 10: More Production Lessons
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
               "More Production Lessons")
add_accent_bar(slide, top=Inches(0.95))

lessons_data = [
    ["Challenge", "What Happened", "Upstream Fix"],
    ["Null Semantics",
     "Velox's collect_set dropped nulls by default; Spark keeps them. TPC-DS queries silently returned wrong results.",
     "Velox #16416, #16947"],
    ["Fallback Storms",
     "One unsupported function forced entire query back to JVM — serialization overhead made native path slower than pure Spark.",
     "Dynamic filter pushdown (Gluten #11657) + Fabric Spark Advisor"],
    ["Security Boundary",
     "Native scan paths must honor Row-Level & Column-Level Security policies.",
     "JNI bridge security enforcement in Fabric NEE"],
    ["Diagnostics",
     "Native segfaults (SIGSEGV) \u2192 no JVM stack trace, opaque crash dumps.",
     "Query trace, Fabric Spark Advisor, native crash symbolication"],
    ["Metrics Gaps",
     "Missing/incomplete metrics for native operators slowed production triage.",
     "Gluten #11861 — improved Generate/join/aggregation metrics fidelity"],
]
add_table(slide, Inches(0.8), Inches(1.2), Inches(11.7), 6, 3, lessons_data,
          col_widths=[Inches(2.2), Inches(5.5), Inches(4.0)], font_size=11)

add_thesis_box(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.7),
    "Key insight:",
    "Production complements rigorous testing with real-world edge cases at previously impossible scale. "
    "Every row above became an upstream PR or a Fabric-specific hardening.")

add_bottom_bar(slide)
add_slide_number(slide, 10)
add_speaker_notes(slide,
    "Each row is a bug or gap discovered at Fabric scale. Null semantics: Velox's collect_set excluded nulls "
    "by default while Spark includes them — TPC-DS silently returned wrong results. Fallback storms: a single "
    "unsupported UDF could force the entire upstream subtree back to JVM — worse than not using Gluten at all. "
    "Security: native scans must honor RLS/CLS policies. Every row became an upstream PR within months.")


# ======================================================================
# SLIDE 11: SECTION DIVIDER — Upstream Contributions
# ======================================================================
section_divider(4, "Upstream Contributions",
                "Validate at scale, upstream fixes, ecosystem benefits", 11)


# ======================================================================
# SLIDE 12: Contribution Themes
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
               "Contribution Themes: Validate \u2192 Upstream \u2192 Ecosystem Benefits")
add_accent_bar(slide, top=Inches(0.95))

add_textbox(slide, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.4),
            "Microsoft contributes to Apache Spark, Apache Gluten, and Velox in collaboration "
            "with Meta, Intel, Kyligence, and the broader community.",
            font_size=12, color=MS_GRAY_90)

contrib_data = [
    ["Theme", "Problem Found in Production", "Upstream Fix (Exemplar PRs)"],
    ["Semantic Parity",
     "Null handling diffs between Spark & Velox broke TPC-DS correctness",
     "Velox: collect_list/collect_set RESPECT NULLS #16328 #16416\n"
     "Gluten: assert_not_null translation #11685\n"
     "Velox: lambda index filter(array,(x,i)\u2192...) #16334"],
    ["Fallback Reduction",
     "Excessive ColumnarToRow transitions from unsupported expressions",
     "Gluten: dynamic filter pushdown #11657\n"
     "Gluten: LOGICAL_PLAN_TAG propagation #11833 (150/150 tests)\n"
     "Velox: setConstantInputs aggregate hook #16624"],
    ["Performance & Scale",
     "Batch-size limits, spill estimation errors, Spark 4.x readiness",
     "Velox: removed FlatNoNulls batch-size limit #16607\n"
     "Velox: PrefixSort memory-estimation overflow fix #15789\n"
     "Gluten: Spark 4.x - 8 test suites enabled #11816"],
    ["Observability",
     "Missing native metrics slowed production triage",
     "Gluten: metrics fidelity for Generate/join/agg #11861\n"
     "Gluten: scan-path formatter reuse #11843"],
    ["TLP Graduation",
     "Post-graduation repo cleanup needed",
     "Gluten: removed incubating refs, CI/release updates #11737-#11742"],
]
add_table(slide, Inches(0.8), Inches(1.6), Inches(11.7), 7, 3, contrib_data,
          col_widths=[Inches(2.0), Inches(4.0), Inches(5.7)], font_size=10)

add_bottom_bar(slide)
add_slide_number(slide, 12)
add_speaker_notes(slide,
    "These are the themes of Microsoft's upstream contributions. Semantic parity was critical: null handling, "
    "lambda indices, array functions — Velox and Spark had subtle differences that broke correctness. "
    "Fallback reduction: we invested in dynamic filter pushdown and logical plan propagation to keep data "
    "columnar. Performance: removed batch-size limits, fixed memory-estimation overflows, enabled Spark 4.x. "
    "The pattern: find it in production, upstream the fix, entire community benefits.")


# ======================================================================
# SLIDE 13: SECTION DIVIDER — Fabric NEE
# ======================================================================
section_divider(5, "Fabric Native Execution Engine",
                "Productionizing Gluten + Velox in Microsoft Fabric", 13)


# ======================================================================
# SLIDE 14: NEE Productionized
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
               "Productionizing Gluten + Velox in Microsoft Fabric")
add_accent_bar(slide, top=Inches(0.95))

# Architecture stack
layers = [
    ("Fabric Runtime 1.3", "Spark 3.5  |  Delta Lake 3.2", MS_BLUE, WHITE),
    ("NEE — Gluten + Velox (C++)", "", MS_PURPLE, WHITE),
    ("OneLake — Delta / Parquet / CSV", "", MS_GREEN, WHITE),
]
for i, (title, sub, color, txt_color) in enumerate(layers):
    y = Inches(1.3 + i * 1.05)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(5.5), Inches(0.75))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    text = f"{title}\n{sub}" if sub else title
    add_textbox(slide, Inches(0.95), y + Inches(0.08), Inches(5.2), Inches(0.6),
                text, font_size=13, color=txt_color, bold=True, alignment=PP_ALIGN.CENTER)
    if i < 2:
        add_textbox(slide, Inches(3.2), y + Inches(0.72), Inches(0.5), Inches(0.3),
                    "\u25bc", font_size=14, color=MS_BLUE, alignment=PP_ALIGN.CENTER)

cfg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(4.55), Inches(3.8), Inches(0.4))
cfg.fill.solid()
cfg.fill.fore_color.rgb = MS_LIGHT_BLUE
cfg.line.fill.background()
add_textbox(slide, Inches(1.5), Inches(4.58), Inches(3.8), Inches(0.35),
            "spark.native.enabled = true", font_size=12,
            color=MS_BLUE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Cascadia Code")

add_textbox(slide, Inches(0.8), Inches(5.1), Inches(5.5), Inches(0.4),
            "Optimal for: complex joins/aggregations, high-volume ETL, "
            "interactive exploration, Delta-heavy workloads",
            font_size=10, color=MS_GRAY_60, alignment=PP_ALIGN.CENTER)

# Metrics with qualifiers — per all 3 reviews
metrics = [
    ("~4\u00d7", "TPC-DS SF1000\n(geometric mean)", "[Internal Testing, Runtime 1.3]", MS_BLUE),
    ("Up to 6\u00d7", "End-to-End\n(query-specific, ETL)", "Ref: GA Blog", MS_GREEN),
    ("~83%", "Runtime Reduction\n(cost savings)", "Ref: Under the Hood Blog", MS_PURPLE),
    ("Included", "In Fabric Runtime\n(no extra cost)", "Ref: No Extra Cost Blog", MS_ORANGE),
]
for i, (val, label, ref, color) in enumerate(metrics):
    row, col = divmod(i, 2)
    x = Inches(6.8 + col * 3.1)
    y = Inches(1.3 + row * 2.0)
    add_metric_card(slide, x, y, Inches(2.8), Inches(1.7), val, label, ref=ref, accent_color=color)

# Methodology footnote — per GPT
add_textbox(slide, Inches(6.8), Inches(5.4), Inches(5.7), Inches(0.3),
            "TPC-DS: synthetic benchmark on homogeneous data. Production results vary by workload.",
            font_size=9, color=MS_GRAY_60, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_slide_number(slide, 14)
add_speaker_notes(slide,
    "This is the product evidence. Fabric Runtime 1.3 is Spark 3.5 + Delta 3.2, with Gluten and Velox built in. "
    "NEE is best for CPU-bound queries: heavy joins, aggregations, ETL. Not every query benefits. "
    "The 4x is geometric mean on TPC-DS SF1000 — synthetic benchmark. Up to 6x on specific ETL queries. "
    "83% is runtime reduction, not billing reduction. Critically: no extra cost for customers. "
    "Be explicit that gains depend on workload shape.")


# ======================================================================
# SLIDE 15: Fabric Extensions
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
               "Fabric Extensions Built on Open-Source Gluten")
add_accent_bar(slide, top=Inches(0.95))

ext_data = [
    ["Category", "Feature", "How It Works"],
    ["Security", "Row-Level Security (RLS)",
     "Policy predicates injected at native scan layer before Velox processes rows — enforced via JNI bridge"],
    ["Security", "Column-Level Security (CLS)",
     "Column masking applied during Substrait plan generation — masked columns never reach Velox in cleartext"],
    ["UDF Support", "Python UDF",
     "Arrow-based columnar transfer via ArrowEvalPython — data stays columnar across JVM-Python, avoiding Row serialization"],
    ["UDF Support", "Scala UDF",
     "JNI-bridged execution with native columnar input — supported UDF classes run without ColumnarToRow transition"],
    ["Performance", "Memory Management",
     "Dynamic reservation/release between JVM & Velox allocators — GC pressure signals trigger native spill before OOM"],
    ["Performance", "Spill-to-Disk",
     "Optimized off-heap spill: sorted spill with compression, async I/O, graceful degradation under pressure"],
    ["Performance", "Hash Probe Optimization",
     "Improved build-side memory layout for cache locality, explicit prefetch hints for probe-side lookups"],
]
add_table(slide, Inches(0.8), Inches(1.2), Inches(11.7), 8, 3, ext_data,
          col_widths=[Inches(1.8), Inches(2.5), Inches(7.4)], font_size=10)

add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.3),
            "These are Fabric-specific production hardening layers built on top of OSS Gluten and Velox.",
            font_size=10, color=MS_GRAY_60)

add_bottom_bar(slide)
add_slide_number(slide, 15)
add_speaker_notes(slide,
    "Fabric extends OSS Gluten with enterprise features. Row-level and column-level security: policies "
    "injected at the native scan layer so Velox never sees data it shouldn't. Python and Scala UDFs: "
    "Arrow-based columnar transfer avoids Row serialization. Memory management: the cooperative "
    "arbitration protocol we discussed earlier. Spill-to-disk with sorted partitions and async I/O. "
    "Hash probe optimization for cache locality. These are Fabric-specific hardening, not upstreamed — "
    "but built on top of OSS Gluten and Velox.")


# ======================================================================
# SLIDE 16: Delta & Observability
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
               "Delta & Storage Optimizations + Observability")
add_accent_bar(slide, top=Inches(0.95))

# LEFT
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.35),
            "Delta & Storage", font_size=18, color=MS_BLUE, bold=True)
delta_data = [
    ["Feature", "Description"],
    ["V-Order", "Default columnar sorting for optimized analytical reads"],
    ["Native Delta Write", "Direct columnar write from Velox to Delta, bypassing JVM row serialization"],
    ["Deletion Vectors", "Native row-level deletes without rewrite"],
    ["Z-Order & Liquid Clustering", "NEE-optimized data layouts for predicate pushdown"],
]
add_table(slide, Inches(0.8), Inches(1.65), Inches(5.5), 5, 2, delta_data,
          col_widths=[Inches(2.0), Inches(3.5)], font_size=11)

# RIGHT
add_textbox(slide, Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.35),
            "Observability", font_size=18, color=MS_BLUE, bold=True)
obs_items = [
    "Spark Advisor: real-time fallback alerts — shows which operator, why, and what to change",
    "Gluten SQL tab: native vs. JVM node visibility in Spark UI",
    "Run Series: before/after NEE performance comparison",
    "In explain() output: Transformer = native C++, ColumnarToRow = back to JVM",
    "Full docs: learn.microsoft.com/fabric/.../native-execution-engine-overview",
]
tb = slide.shapes.add_textbox(Inches(6.8), Inches(1.65), Inches(5.7), Inches(3.0))
tf = tb.text_frame
tf.word_wrap = True
add_bullet_list(tf, obs_items, font_size=12)

add_textbox(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.5),
            "Refs: learn.microsoft.com/fabric/data-engineering/spark-advisor-introduction  |  "
            "learn.microsoft.com/fabric/data-engineering/delta-optimization-and-v-order",
            font_size=9, color=MS_GRAY_60)

add_bottom_bar(slide)
add_slide_number(slide, 16)
add_speaker_notes(slide,
    "Fabric added three key observability tools. Spark Advisor shows which operators fall back to JVM "
    "and why — no need to manually parse explain() output. Gluten SQL tab in Spark UI shows native vs. "
    "JVM nodes side-by-side. Run Series compares before/after NEE performance. V-Order is Fabric's "
    "default columnar sort. Native Delta writes bypass Row serialization. These complete the "
    "production-readiness picture.")


# ======================================================================
# SLIDE 17: SECTION DIVIDER — Graduation & Roadmap
# ======================================================================
section_divider(6, "Graduation & Roadmap",
                "From incubation to TLP, reinforced by production adoption and multi-company stewardship", 17)


# ======================================================================
# SLIDE 18: TLP Journey & Lessons
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
               "Gluten's TLP Journey & Lessons")
add_accent_bar(slide, top=Inches(0.95))

# Timeline
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.35),
            "Timeline", font_size=18, color=MS_BLUE, bold=True)
timeline = [
    ("2022 Q1", "Intel & Kyligence create Gluten"),
    ("Dec 2023", "Entered Apache Incubator"),
    ("2024\u201325", "Microsoft adoption \u2014 Fabric NEE in production, upstream contributions begin"),
    ("Jan 2025", "v1.3.0 \u2014 OOM \u2193 70%"),
    ("Feb 18, 2026", "Apache TLP \u2014 1,500+ stars, ~194 contributors"),
]
for i, (date, desc) in enumerate(timeline):
    y = Inches(1.7 + i * 0.65)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.05), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = MS_BLUE if i < 4 else MS_GREEN
    dot.line.fill.background()
    if i < len(timeline) - 1:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.06), y + Inches(0.2), Inches(0.03), Inches(0.5))
        line.fill.solid()
        line.fill.fore_color.rgb = MS_GRAY_30
        line.line.fill.background()
    add_textbox(slide, Inches(1.3), y - Inches(0.02), Inches(1.5), Inches(0.3),
                date, font_size=11, color=MS_BLUE, bold=True)
    add_textbox(slide, Inches(2.8), y - Inches(0.02), Inches(3.2), Inches(0.5),
                desc, font_size=11, color=MS_GRAY_90)

# Lessons
add_textbox(slide, Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.35),
            "Key Lessons", font_size=18, color=MS_BLUE, bold=True)
add_card(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(1.0),
         "Governance Matters as Much as Code",
         "ASF graduation required a PMC with Intel, Kyligence, Microsoft, and independents "
         "\u2014 not a single company's project.",
         accent_color=MS_BLUE)
add_card(slide, Inches(6.8), Inches(2.9), Inches(5.5), Inches(1.0),
         "Memory: The Hardest Problem",
         "JVM vs. native off-heap \u2014 9-month collaboration to build cooperative memory arbitration. OOM dropped 70%.",
         accent_color=MS_ORANGE)
add_card(slide, Inches(6.8), Inches(4.1), Inches(5.5), Inches(1.2),
         "Production \u2192 Upstream \u2192 Graduation",
         "How production adoption (at Microsoft and elsewhere) validated Gluten at scale "
         "and accelerated TLP readiness \u2014 not a downstream fork strategy.",
         accent_color=MS_GREEN)

add_textbox(slide, Inches(0.8), Inches(5.5), Inches(5.5), Inches(0.25),
            "Ref: incubator.apache.org/projects/gluten.html", font_size=9, color=MS_GRAY_60)

add_bottom_bar(slide)
add_slide_number(slide, 18)
add_speaker_notes(slide,
    "Intel and Kyligence created Gluten in 2022. Microsoft adopted it for Fabric NEE in 2024-25, "
    "started upstream contributions. Key lessons: Governance comes before code \u2014 ASF graduation "
    "required a diverse PMC. Memory was the hardest problem \u2014 9-month collaboration. "
    "Our strategy was never to fork; it was to find problems at scale and upstream fixes.")


# ======================================================================
# SLIDE 19: 2026 Roadmap
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
               "2026 Roadmap & Microsoft Focus")
add_accent_bar(slide, top=Inches(0.95))

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.35),
            "Gluten Community Roadmap", font_size=18, color=MS_BLUE, bold=True)
roadmap_data = [
    ["Q", "Version", "Highlights"],
    ["Q1", "1.6.0", "Post-grad cleanup, Spark 4.0 beta"],
    ["Q2", "1.7.0", "Multi-core, PySpark UDF"],
    ["Q3", "1.8.0", "Flink PoC, GPU exploration"],
    ["Q4", "2.0.0", "Spark 4.0 GA, new backends"],
]
add_table(slide, Inches(0.8), Inches(1.7), Inches(5.5), 5, 3, roadmap_data,
          col_widths=[Inches(0.6), Inches(1.0), Inches(3.9)], font_size=12)
add_textbox(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.25),
            "Ref: github.com/apache/gluten/issues/11827", font_size=9, color=MS_GRAY_60)

add_textbox(slide, Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.35),
            "Microsoft Focus Areas", font_size=18, color=MS_BLUE, bold=True)
add_card(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(0.9),
         "Spark Compatibility",
         "Spark 4.x semantic parity, ANSI mode, TIMESTAMP_NTZ, function coverage",
         accent_color=MS_BLUE)
add_card(slide, Inches(6.8), Inches(2.8), Inches(5.5), Inches(0.9),
         "Performance",
         "Dynamic filters, bloom-filter pushdown, join optimization, memory & spill robustness",
         accent_color=MS_GREEN)
add_card(slide, Inches(6.8), Inches(3.9), Inches(5.5), Inches(0.9),
         "Enterprise Features",
         "Structured Streaming, complex types (array/map/struct), semi-structured data (JSON variant shredding)",
         accent_color=MS_PURPLE)

add_thesis_box(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(0.8),
    "Shared investment:",
    "Each new Spark-compatible Velox capability expands native execution opportunities "
    "for Fabric and the broader ecosystem. We're committed to closing the Spark-Velox "
    "semantic gap \u2014 every function we upstream benefits all adopters.")

add_bottom_bar(slide)
add_slide_number(slide, 19)
add_speaker_notes(slide,
    "Community roadmap: Q1 cleanup, Q2 multi-core and PySpark UDFs, Q3 Flink PoC and GPU exploration, "
    "Q4 Spark 4.0 GA. Microsoft's focus: Spark 4.x parity, ANSI mode, performance (dynamic filters, "
    "bloom-filter pushdown), and enterprise features (Structured Streaming, complex types). "
    "Every function we upstream benefits the entire community. "
    "If you're using Spark 4.0 or Structured Streaming, join us.")


# ======================================================================
# SLIDE 20: SECTION DIVIDER — Call for Action
# ======================================================================
section_divider(7, "Call for Action",
                "We welcome deeper collaboration with the Velox and Gluten communities.", 20)


# ======================================================================
# SLIDE 21: Join the Community
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
               "Join the Community")
add_accent_bar(slide, top=Inches(0.95))

add_textbox(slide, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.35),
            "Gluten and Velox thrive on diverse contributors. Here's how to get involved:",
            font_size=14, color=MS_GRAY_90)

cta_items = [
    ("Contribute Code",
     "Start here: add a missing Spark function to Velox (e.g., approx_percentile, avg_if). "
     "It unblocks native execution for thousands of users. "
     "See help-wanted issues on Gluten and Velox GitHub.",
     MS_BLUE),
    ("Test & Report",
     "Run Gluten on your workloads. File issues with explain() output and "
     "Velox error messages \u2014 production edge cases are the most valuable contributions.",
     MS_GREEN),
    ("Improve Documentation",
     "Configuration guides, fallback analysis tutorials, and migration recipes are always needed. "
     "Docs PRs welcome at apache/gluten.",
     MS_PURPLE),
    ("Integrate & Extend",
     "Build new backends, connect to your storage layer, or prototype Flink/Streaming integrations. "
     "The plugin architecture is designed for extensibility.",
     MS_ORANGE),
]
for i, (title, body, color) in enumerate(cta_items):
    row, col = divmod(i, 2)
    x = Inches(0.8 + col * 6.0)
    y = Inches(1.7 + row * 1.7)
    add_card(slide, x, y, Inches(5.7), Inches(1.4), title, body, accent_color=color)

add_thesis_box(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.7),
    "Open invitation:",
    "If you are debugging difficult Spark-on-Velox edge cases, collaborate with us upstream "
    "so the fixes benefit the full ecosystem.")

add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.4),
            "github.com/apache/gluten  |  github.com/facebookincubator/velox  |  "
            "learn.microsoft.com/fabric/.../native-execution-engine-overview",
            font_size=10, color=MS_GRAY_60, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_slide_number(slide, 21)
add_speaker_notes(slide,
    "Four ways to contribute. Code: add a missing Spark function to Velox — "
    "approx_percentile, avg_if, etc. It unblocks execution for thousands of users. "
    "Testing: run Gluten on your workloads, file issues with explain() output. "
    "Docs: configuration guides, migration recipes. Integrate: new backends, "
    "Flink prototypes, storage connectors. We want your hardest edge cases in the open.")


# ======================================================================
# SLIDE 22: CLOSING
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, MS_BLUE, MS_DARK_BLUE)
add_ms_logo(slide, Inches(5.9), Inches(0.8), Inches(0.55))
add_textbox(slide, Inches(6.55), Inches(0.82), Inches(2), Inches(0.55),
            "Microsoft", font_size=22, color=WHITE, bold=True)

add_textbox(slide, Inches(1.5), Inches(1.7), Inches(10.3), Inches(0.6),
            "Thank You & Q&A",
            font_size=36, color=WHITE, bold=True,
            font_name="Segoe UI Semibold", alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.4), Inches(10.3), Inches(0.6),
            "From Apache Incubator to Top-Level Project:\n"
            "The Apache Gluten Journey with Velox at Microsoft",
            font_size=18, color=RGBColor(0xB4, 0xD6, 0xFA), alignment=PP_ALIGN.CENTER)

# Speakers
add_textbox(slide, Inches(2.5), Inches(3.3), Inches(8.3), Inches(0.35),
            "William Chen  |  Apache Gluten PMC Member  |  Microsoft",
            font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.5), Inches(3.65), Inches(8.3), Inches(0.3),
            "williamchen@microsoft.com",
            font_size=12, color=RGBColor(0xB4, 0xD6, 0xFA), alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.5), Inches(4.1), Inches(8.3), Inches(0.35),
            "Ajith Shetty  |  Microsoft",
            font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.5), Inches(4.45), Inches(8.3), Inches(0.3),
            "ajithshetty@microsoft.com",
            font_size=12, color=RGBColor(0xB4, 0xD6, 0xFA), alignment=PP_ALIGN.CENTER)

# Links
links = [
    "github.com/apache/gluten",
    "github.com/facebookincubator/velox",
    "learn.microsoft.com/fabric/data-engineering/native-execution-engine-overview",
]
for i, link in enumerate(links):
    add_textbox(slide, Inches(2.5), Inches(5.1 + i * 0.3), Inches(8.3), Inches(0.3),
                link, font_size=12, color=RGBColor(0x8E, 0xC8, 0xFF), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(6.1), Inches(10.3), Inches(0.35),
            "Shared commitment: upstream innovation, production hardening at scale, "
            "and shared progress across the Velox ecosystem.",
            font_size=13, color=RGBColor(0xB4, 0xD6, 0xFA), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(6.8), Inches(10.3), Inches(0.3),
            "VeloxCon 2026  |  April 30  |  Meta HQ, Menlo Park",
            font_size=10, color=RGBColor(0x8E, 0xC8, 0xFF), alignment=PP_ALIGN.CENTER)


# ======================================================================
# SAVE
# ======================================================================
output_path = r"C:\Users\williamchen\Projects\presentations\presentations\veloxcon2026-gluten-microsoft\VeloxCon2026_Gluten_Microsoft.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
